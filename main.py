!pip install python-docx python-pptx mammoth pytesseract pymupdf chromadb sentence-transformers
#pdfplumber

from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import pandas as pd
from docx import Document as DocxDocument
import mammoth
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET
import requests
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image # scanned pdf extractor
import fitz # PDF new extractor
import re
import os
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
from openai import OpenAI

BASE_DATA_DIR = Path("/content/data")  # where you’ll put your files in Colab
BASE_DATA_DIR.mkdir(parents=True, exist_ok=True)

VECTOR_DB_DIR = "/content/chroma_db"
# Initialize Chroma client
chroma_client = chromadb.PersistentClient(
    path=VECTOR_DB_DIR,
    settings=Settings(anonymized_telemetry=False)
)
COLLECTION_NAME = "arabic_rag_collection"
EMBEDDING_MODEL_NAME = "sentence-transformers/paraphrase-multilingual-mpnet-base-v2"

# Initialize embedding model (CPU)
embedding_model = SentenceTransformer(EMBEDDING_MODEL_NAME)
# Set your keys here or via environment variables in Colab
OPENAI_API_KEY = "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"
TELEGRAM_BOT_TOKEN = "8527483344:AAHg69NPJie_ZHJCrs2bGVoGDCtl7tf84iQ"  # <<< change this

os.environ["OPENAI_API_KEY"] = OPENAI_API_KEY

client = OpenAI(api_key=OPENAI_API_KEY)

def load_pdf(path: Path) -> str:
    doc = fitz.open(path)
    all_text = []

    for page_index in range(doc.page_count):
        # Load ONLY the specific page into RAM
        page = doc.load_page(page_index)

        # Extract text (very lightweight)
        text = page.get_text("text", flags=fitz.TEXT_PRESERVE_LIGATURES)

        # If empty fallback to block extraction (less RAM than OCR)
        if not text.strip():
            blocks = page.get_text("blocks")
            text = " ".join(b[4] for b in blocks if len(b) >= 5)

        all_text.append(text)

        # Release memory for the loaded page
        del page

    doc.close()

    return "\n".join(all_text)

def load_docx(path: Path) -> str: # issue in the definition
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
      with open(path, "rb") as f:
        result = mammoth.convert_to_html(f)
        html = result.value

      # Now extract text from HTML
      soup = BeautifulSoup(html, "html.parser")

      paragraphs = soup.get_text(" ", strip=True)
      return paragraphs
    return "\n".join(paragraphs)


def load_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slide_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
    if not slide_texts:
      texts = []
      # Open PPTX as a zip archive
      with zipfile.ZipFile(path, 'r') as pptx:
        # Iterate through all slide files
        for slide_name in pptx.namelist():
            if slide_name.startswith("ppt/slides/slide") and slide_name.endswith(".xml"):
                xml_content = pptx.read(slide_name)
                # Parse XML using builtin xml.etree (NO LXML)
                root = ET.fromstring(xml_content)
                # PPTX text nodes are usually inside <a:t> tags
                namespace = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                for node in root.iter():
                    if node.tag.endswith("}t"):  # extract <a:t> text
                        if node.text:
                            texts.append(node.text.strip())
      return "\n".join(texts)
    return "\n".join(slide_texts)

def load_excel_or_csv(path: Path) -> str:
    # Excel or CSV to text: concatenate all cell values
    if path.suffix.lower() == ".csv":
        df = pd.read_csv(path)
    else:
        df = pd.read_excel(path)

    # Convert all cells to strings, drop NaN
    df = df.astype(str)
    # You can customize how you flatten the table
    return "\n".join(df.apply(lambda row: " | ".join(row.values), axis=1).tolist())

def load_website(url: str) -> str:
    resp = requests.get(url, timeout=10)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "lxml")

    # Remove script/style
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()

    texts = []
    for element in soup.find_all(["h1", "h2", "h3", "h4", "p", "li"]):
        text = element.get_text(separator=" ", strip=True)
        if text:
            texts.append(text)
    return "\n".join(texts)


def load_any(source: str) -> Tuple[str, Dict[str, Any]]:
    """
    source: file path (str) or URL (starting with http)
    Returns: (text, metadata)
    """
    if source.startswith("http://") or source.startswith("https://"):
        text = load_website(source)
        metadata = {
            "type": "website",
            "source": source,
            "title": source  # you can later parse <title> if you want
        }
        return text, metadata

    path = Path(source)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        text = load_pdf(path)
        doc_type = "pdf"
    elif suffix in [".docx"]:
        text = load_docx(path)
        doc_type = "docx"
    elif suffix in [".pptx"]:
        text = load_pptx(path)
        doc_type = "pptx"
    elif suffix in [".xls", ".xlsx", ".csv"]:
        text = load_excel_or_csv(path)
        doc_type = "spreadsheet"
    else:
        # Fallback: try to read as plain text
        text = Path(path).read_text(encoding="utf-8", errors="ignore")
        doc_type = "text"

    metadata = {
        "type": doc_type,
        "source": str(path),
        "title": path.name,
    }
    return text, metadata

sample_path = "/content/sample_data/ArrivalPost-Arb.pdf"  # change this
text, meta = load_any(sample_path)
print(meta)
print(text[:1000])

ARABIC_DIACRITICS = re.compile(r"[\u0610-\u061A\u064B-\u065F\u06D6-\u06DC\u06DF-\u06E8\u06EA-\u06ED]")

def normalize_arabic(text: str) -> str:
    # Remove tatweel
    text = text.replace("ـ", "")

    # Normalize alef forms
    text = re.sub("[أإآا]", "ا", text)
    # Normalize taa marbuta to haa? (optional – can change search semantics)
    text = text.replace("ة", "ه")

    # Remove diacritics
    text = ARABIC_DIACRITICS.sub("", text)
    # Remove markdown bold/italic
    text = re.sub(r"\*+", "", text)

    # Remove bullet prefixes
    text = re.sub(r"^[\-\*\•]+\s*", "", text, flags=re.MULTILINE)

    # Normalize spaces
    text = re.sub(r"\s+", " ", text).strip()

    return text


def clean_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\u200f", "")  # RTL mark
    text = text.replace("\u200e", "")  # LTR mark
    text = text.replace("\ufeff", "")  # BOM
    return text

def chunk_text(text, max_chars=800, overlap=200):
    chunks = []
    buffer = ""
    overlap_buffer = ""

    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue

        # Add line to current buffer
        if len(buffer) + len(line) + 1 <= max_chars:
            buffer += " " + line
        else:
            # Finalize chunk
            chunk = buffer.strip()
            if chunk:
                chunks.append(chunk)

            # Start new buffer with overlap
            if len(buffer) > overlap:
                overlap_buffer = buffer[-overlap:]
            else:
                overlap_buffer = buffer

            buffer = overlap_buffer + " " + line

    # Add last chunk
    if buffer.strip():
        chunks.append(buffer.strip())

    return chunks

chunks = chunk_text(text)
print("Number of chunks:", len(chunks))
print(chunks[0][:400])

collection = chroma_client.get_or_create_collection(
    name=COLLECTION_NAME,
    metadata={"hnsw:space": "cosine"},
)

def embed_texts(texts: List[str]) -> List[List[float]]:
    # SentenceTransformer returns numpy array; convert to list
    embeddings = embedding_model.encode(texts, batch_size=32, show_progress_bar=True)
    return [emb.tolist() for emb in embeddings]


import uuid
from tqdm.auto import tqdm

def index_source(source: str, doc_id_prefix: Optional[str] = None):
    text, base_metadata = load_any(source)
    chunks = chunk_text(text)

    if not chunks:
        print(f"No valid chunks for source: {source}")
        return

    print(f"Indexing {len(chunks)} chunks from {source}")

    embeddings = embed_texts(chunks)

    ids = []
    metadatas = []
    for i, chunk in enumerate(chunks):
        chunk_id = f"{doc_id_prefix or str(uuid.uuid4())}_{i}"
        ids.append(chunk_id)
        md = base_metadata.copy()
        md.update({
            "chunk_index": i,
        })
        metadatas.append(md)

    collection.add(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    print("Done indexing.")

def index_many(sources: List[str]):
    for src in tqdm(sources):
        try:
            index_source(src)
        except Exception as e:
            print(f"Error indexing {src}: {e}")

sources = [
    "/content/sample_data/ArrivalPost-Arb.pdf",
    "https://www.moh.gov.sa/HealthAwareness/EducationalContent/wh/Pages/Hypothyroidism-and-Pregnancy.aspx",
    "https://www.moh.gov.sa/HealthAwareness/EducationalContent/wh/Pages/001.aspx",
    "/content/sample_data/2018-11-13-002.pdf"
]
index_many(sources)

def retrieve_relevant_chunks(
    query: str,
    top_k: int = 5,
) -> List[Dict[str, Any]]:
    normalized_query = normalize_arabic(query)
    query_embedding = embed_texts([normalized_query])[0]

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=top_k,
        include=["documents", "metadatas", "distances"],
    )

    chunks = []
    for doc, meta, dist in zip(
        results["documents"][0],
        results["metadatas"][0],
        results["distances"][0],
    ):
        chunks.append({
            "text": doc,
            "metadata": meta,
            "score": float(dist),
        })
    return chunks

def clean_context_text(text: str) -> str:
    """
    Cleans chunk text before sending to the LLM:
    - Removes Markdown (**bold**, *italic*)
    - Removes bullets
    - Removes repeated punctuation
    - Normalizes spacing
    - Keeps readable Arabic-only content
    """

    # Remove markdown bold/italic markers
    text = re.sub(r"\*{1,3}", "", text)

    # Remove markdown-like numbered items like **1. نص**
    text = re.sub(r"\*+\s*(\d+)\s*\.*\s*\*+", r"\1.", text)

    # Remove leading bullets
    text = re.sub(r"^[\-\*\•]+\s*", "", text, flags=re.MULTILINE)

    # Remove duplicated punctuation
    text = re.sub(r"([\.؟!])\1+", r"\1", text)

    # Remove multiple spaces
    text = re.sub(r"\s+", " ", text)

    return text.strip()

def build_prompt_arabic(question: str, contexts: List[Dict[str, Any]]) -> str:
    """
    Builds a clean Arabic prompt with:
    - Cleaned contextual chunks
    - No Markdown
    - Strict source citation rules
    """
    context_lines = []

    for i, c in enumerate(contexts, start=1):
        src = c["metadata"].get("title") or c["metadata"].get("source") or f"chunk_{i}"
        clean_text = clean_context_text(c["text"])

        context_lines.append(
            f"[المصدر {i}: {src}]\n{clean_text}"
        )

    context_block = "\n\n".join(context_lines)

    prompt = f"""
السؤال: {question}

النصوص المسترجعة من المصادر:
{context_block}

التعليمات:
- أجب إجابة واضحة ومنسقة باللغة العربية فقط.
- لا تستخدم أي علامات Markdown مثل ** أو * أو -.
- لا تستخدم ترقيم تلقائي غير متسق.
- استخدم فقط المعلومات الموجودة في النصوص أعلاه.
- لا تضف أي مصادر غير الموجودة في القائمة أعلاه.
- في نهاية الإجابة، لا تذكر أي مصادر. سيتم إظهارها تلقائياً لاحقاً.

الإجابة:
"""

    return prompt.strip()

def generate_answer_from_llm(prompt: str, model: str = "gpt-4o-mini") -> str:
    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": (
                "أنت مساعد ذكي متخصص في الإجابة باللغة العربية بشكل واضح ودقيق. "
                "لا تستخدم Markdown. لا تستخدم **. "
                "لا تضف أي مصادر إضافية غير التي تظهر في المقاطع."
            )},
            {"role": "user", "content": prompt}
        ],
        max_tokens=800,
        temperature=0.1
    )

    answer = response.choices[0].message.content

    # Remove any hallucinated "المصادر المستخدمة"
    answer = re.split(r"المصادر المستخدمة\s*[:：]?", answer)[0].strip()

    return answer

def answer_question(question: str, top_k: int = 5) -> Tuple[str, List[Dict[str, Any]]]:
    # Retrieve
    chunks = retrieve_relevant_chunks(question, top_k=top_k)

    if not chunks:
        return "لم أستطع إيجاد أي معلومات مرتبطة بسؤالك في قاعدة المعرفة.", []

    # Build prompt
    prompt = build_prompt_arabic(question, chunks)

    # Generate
    answer = generate_answer_from_llm(prompt)

    return answer, chunks

def extract_used_source_numbers(answer: str):
    """
    Extracts 'المصدر X' references from the model’s answer.
    Returns a set of source numbers actually used.
    """
    matches = re.findall(r"\(المصدر\s+(\d+)\)", answer)
    return set(int(m) for m in matches)

q = "ما هي أعراض ألم ما بعد الولادة الطبيعية؟"

ans, retrieved_chunks = answer_question(q, top_k=5)

print("الإجابة:\n")
print(ans)
used_numbers = extract_used_source_numbers(ans)

print("\nالمصادر المستخدمة:")
already_listed = set()

unique_sources = set()

# CASE 1: Model cited sources
if used_numbers:
    for idx in sorted(used_numbers):
        if 1 <= idx <= len(retrieved_chunks):
            c = retrieved_chunks[idx - 1]
            src = c['metadata'].get('title') or c['metadata'].get('source')
            if src not in unique_sources:
                print(f"- {src}")
                unique_sources.add(src)

# CASE 2: Model did NOT cite sources → fallback to unique retrieved sources
else:
    for c in retrieved_chunks:
        src = c['metadata'].get('title') or c['metadata'].get('source')
        if src not in unique_sources:
            print(f"- {src}")
            unique_sources.add(src)
